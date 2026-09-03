"""wiki-parse parser 单测。覆盖分派：MD/TXT/HTML + PDF 缺失分支 +
DOCX/XLSX/PPTX/EPUB（样本文件由对应库在测试内现造，不提交二进制 fixture）。

真 PDF 文本提取需 fixture（reportlab 构造或 sample.pdf），MVP 跳，排期补。
"""

from __future__ import annotations

import io
import sys
import zipfile

import pytest

from wiki_parse.parser import ParseError, extract


def test_extract_markdown_passthrough():
    data = b"# Title\n\nbody text\n- item"
    assert extract(data, filename="note.md") == "# Title\n\nbody text\n- item"


def test_extract_txt_utf8_replace_on_invalid_bytes():
    # 非法 utf-8 字节 → errors=replace 不抛，保留可解码部分
    data = b"hello\xff\xffworld"
    out = extract(data, filename="log.txt")
    assert "hello" in out and "world" in out


def test_extract_html_strips_script_style_title():
    html = (
        b"<html><head><title>IGNORE</title>"
        b"<style>body{}</style></head>"
        b"<body><p>Hello <b>world</b></p>"
        b"<script>alert(1)</script></body></html>"
    )
    out = extract(html, filename="page.html")
    assert "Hello" in out and "world" in out
    assert "alert" not in out     # script stripped
    assert "IGNORE" not in out    # title stripped
    assert "body{}" not in out    # style stripped


def test_extract_html_by_content_sniff_without_extension():
    # 无 .html 扩展名、mime 空 —— 靠 doctype 探测
    data = b"<!doctype html><html><body><p>sniffed</p></body></html>"
    out = extract(data, filename="unknown")
    assert "sniffed" in out


def test_extract_html_boilerplate_only_raises():
    data = b"<html><head><title>x</title></head><body></body></html>"
    with pytest.raises(ParseError):
        extract(data, filename="empty.html")


def test_extract_empty_file_raises():
    with pytest.raises(ParseError, match="empty"):
        extract(b"", filename="x.txt")


def test_extract_pdf_missing_pypdf_raises(monkeypatch):
    # 模拟 pypdf 未安装：sys.modules["pypdf"] = None 让 import 抛 ModuleNotFoundError
    monkeypatch.setitem(sys.modules, "pypdf", None)
    with pytest.raises(ParseError, match="pypdf"):
        extract(b"%PDF-1.4 fake", mime="application/pdf", filename="x.pdf")


def test_extract_decoded_empty_raises():
    with pytest.raises(ParseError, match="decoded text empty"):
        extract(b"   \n\t  ", filename="blank.txt")


# ─── Office / EPUB 格式（样本由对应库现造）─────────────────────────────


def _make_docx(paragraphs: list[str]) -> bytes:
    """最小合法 DOCX（只 content-types + rels + document.xml）。"""
    body = "".join(
        f"<w:p><w:r><w:t>{p}</w:t></w:r></w:p>" for p in paragraphs
    )
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        z.writestr(
            "[Content_Types].xml",
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/>'
            '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
            "</Types>",
        )
        z.writestr(
            "_rels/.rels",
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>'
            "</Relationships>",
        )
        z.writestr(
            "word/document.xml",
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
            f"<w:body>{body}</w:body></w:document>",
        )
    return buf.getvalue()


def _make_xlsx(sheets: dict[str, list[list[str]]]) -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    wb.remove(wb.active)
    for title, rows in sheets.items():
        ws = wb.create_sheet(title=title)
        for row in rows:
            ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_pptx(slide_texts: list[list[str]]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    for texts in slide_texts:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        box = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(4), Inches(2)
        )
        box.text_frame.text = "\n".join(texts)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _make_epub(chapters: dict[str, str]) -> bytes:
    from ebooklib import epub

    book = epub.EpubBook()
    book.set_identifier("test-book")
    book.set_title("Test Book")
    book.set_language("en")
    items = []
    for name, html in chapters.items():
        c = epub.EpubHtml(title=name, file_name=f"{name}.xhtml", lang="en")
        c.content = html
        book.add_item(c)
        items.append(c)
    book.toc = items
    book.spine = ["nav", *items]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    buf = io.BytesIO()
    epub.write_epub(buf, book)
    return buf.getvalue()


def test_extract_docx_by_extension():
    data = _make_docx(["Hello DOCX world", "second paragraph"])
    out = extract(data, filename="report.docx")
    assert "Hello DOCX world" in out
    assert "second paragraph" in out


def test_extract_docx_by_mime_without_extension():
    data = _make_docx(["mime dispatch works"])
    out = extract(
        data,
        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        filename="blob",
    )
    assert "mime dispatch works" in out


def test_extract_docx_empty_raises():
    with pytest.raises(ParseError, match="DOCX"):
        extract(_make_docx([]), filename="empty.docx")


def test_extract_docx_missing_mammoth_raises(monkeypatch):
    monkeypatch.setitem(sys.modules, "mammoth", None)
    with pytest.raises(ParseError, match="mammoth"):
        extract(b"fake", filename="x.docx")


def test_extract_xlsx_multiple_sheets():
    data = _make_xlsx(
        {
            "Summary": [["name", "score"], ["alpha", "42"]],
            "Empty": [],
            "Detail": [["only cell"]],
        }
    )
    out = extract(data, filename="book.xlsx")
    assert "--- sheet: Summary ---" in out
    assert "name | score" in out
    assert "alpha | 42" in out
    assert "--- sheet: Detail ---" in out
    assert "Empty" not in out  # 空 sheet 不产出段落


def test_extract_xlsx_empty_raises():
    with pytest.raises(ParseError, match="XLSX"):
        extract(_make_xlsx({"S": []}), filename="empty.xlsx")


def test_extract_xlsx_missing_openpyxl_raises(monkeypatch):
    monkeypatch.setitem(sys.modules, "openpyxl", None)
    with pytest.raises(ParseError, match="openpyxl"):
        extract(b"fake", filename="x.xlsx")


def test_extract_pptx_per_slide():
    data = _make_pptx([["Title one", "bullet A"], ["Second slide text"]])
    out = extract(data, filename="deck.pptx")
    assert "--- slide 1 ---" in out
    assert "Title one" in out and "bullet A" in out
    assert "--- slide 2 ---" in out
    assert "Second slide text" in out


def test_extract_pptx_empty_raises():
    with pytest.raises(ParseError, match="PPTX"):
        extract(_make_pptx([]), filename="empty.pptx")


def test_extract_pptx_missing_python_pptx_raises(monkeypatch):
    monkeypatch.setitem(sys.modules, "pptx", None)
    with pytest.raises(ParseError, match="python-pptx"):
        extract(b"fake", filename="x.pptx")


def test_extract_epub_chapters_stripped():
    data = _make_epub(
        {
            "ch1": "<html><body><h1>Chapter One</h1>"
                   "<p>EPUB body text here</p>"
                   "<script>ignore_me()</script></body></html>",
            "ch2": "<html><body><p>second chapter content</p></body></html>",
        }
    )
    out = extract(data, filename="book.epub")
    assert "EPUB body text here" in out
    assert "second chapter content" in out
    assert "ignore_me" not in out  # script 被 tag-strip 掉


def test_extract_epub_by_mime_without_extension():
    data = _make_epub({"c": "<html><body><p>epub mime dispatch</p></body></html>"})
    out = extract(data, mime="application/epub+zip", filename="blob")
    assert "epub mime dispatch" in out


def test_extract_epub_missing_ebooklib_raises(monkeypatch):
    monkeypatch.setitem(sys.modules, "ebooklib", None)
    with pytest.raises(ParseError, match="ebooklib"):
        extract(b"fake", filename="x.epub")


def test_extract_epub_empty_raises():
    # 唯一章节只有 script（tag-strip 后无正文；nav 目录被跳过）→ ParseError。
    # 注意不能写真空 body —— ebooklib write_epub 对空文档自身就抛 lxml ParserError。
    data = _make_epub({"c": "<html><body><script>only_script()</script></body></html>"})
    with pytest.raises(ParseError, match="EPUB"):
        extract(data, filename="empty.epub")


# ─── PDF（最小合法 PDF 手工拼 xref，不提交二进制 fixture）────────────────


def _make_pdf(page_texts: list[str | None]) -> bytes:
    """构造最小合法 PDF。page_texts[i] 为 None → 该页无文本层（空页）。

    文本页用 Helvetica Type1 + 简单 Tj content stream，pypdf 可 extract_text。
    """
    n = len(page_texts)
    font_id = 3
    page_ids = [4 + i for i in range(n)]
    content_ids: dict[int, int] = {}
    next_id = 4 + n
    for i, t in enumerate(page_texts):
        if t is not None:
            content_ids[i] = next_id
            next_id += 1

    bodies: dict[int, bytes] = {
        1: b"<< /Type /Catalog /Pages 2 0 R >>",
        2: (
            "<< /Type /Pages /Kids ["
            + " ".join(f"{pid} 0 R" for pid in page_ids)
            + f"] /Count {n} >>"
        ).encode(),
        font_id: b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    }
    for i, t in enumerate(page_texts):
        if t is None:
            bodies[page_ids[i]] = (
                b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] >>"
            )
        else:
            bodies[page_ids[i]] = (
                f"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
                f"/Resources << /Font << /F1 {font_id} 0 R >> >> "
                f"/Contents {content_ids[i]} 0 R >>"
            ).encode()
            stream = f"BT /F1 12 Tf 72 720 Td ({t}) Tj ET".encode()
            bodies[content_ids[i]] = (
                f"<< /Length {len(stream)} >>\nstream\n".encode()
                + stream + b"\nendstream"
            )

    out = io.BytesIO()
    out.write(b"%PDF-1.4\n")
    offsets = {}
    for oid in range(1, next_id):
        offsets[oid] = out.tell()
        out.write(f"{oid} 0 obj\n".encode() + bodies[oid] + b"\nendobj\n")
    xref_pos = out.tell()
    out.write(f"xref\n0 {next_id}\n".encode())
    out.write(b"0000000000 65535 f \n")
    for oid in range(1, next_id):
        out.write(f"{offsets[oid]:010d} 00000 n \n".encode())
    out.write(
        f"trailer\n<< /Size {next_id} /Root 1 0 R >>\n"
        f"startxref\n{xref_pos}\n%%EOF\n".encode()
    )
    return out.getvalue()


def test_extract_pdf_empty_page_gets_placeholder():
    # 混合 PDF：第 1 页有文本层，第 2 页空 —— 空页不静默丢，落占位符
    data = _make_pdf(["Hello PDF world", None])
    out = extract(data, filename="mix.pdf")
    assert "--- page 1 ---" in out
    assert "Hello PDF world" in out
    assert "[page 2: 无文本层]" in out


def test_extract_pdf_all_blank_raises_pointing_to_ocr():
    data = _make_pdf([None, None])
    with pytest.raises(ParseError, match="OCR") as exc_info:
        extract(data, filename="scan.pdf")
    # 文案指向 OCR 未启用/失败，不再是旧的"OCR 排期"
    assert "OCR 排期" not in str(exc_info.value)
