"""File → extracted text. Parsers: PDF (pypdf) + DOCX (mammoth，对齐
docproc-web 客户端侧 mammoth.js) + XLSX (openpyxl) + PPTX (python-pptx)
+ EPUB (ebooklib + stdlib tag-strip) + MD/TXT (utf-8) + HTML (stdlib tag-strip)。

分派按 mime + filename 扩展名 + 内容探测。失败抛 ``ParseError``，runner
捕获后回写 ``parse_status='error'``。

边界（llm_wiki 有而我们暂未做的）：
- HTML 只 stdlib strip tag，无 readability/trafilatura 正文抽取
- PDF 只 pypdf 文本层提取，扫描版/复杂表格效果差 —— OCR 由 runner 经
  ocr.py 调自部署 MinerU 补齐（B1，OCR 启用时全量 PDF 走 MinerU）
- 无 MOBI（排期）

content_hash = sha256(extracted_text)（不是文件字节），所以同内容不同格式
会被判重 —— 这是产品决策（用户已确认）。
"""

from __future__ import annotations

import io
from html.parser import HTMLParser
from typing import List, Optional


class ParseError(RuntimeError):
    """文件解析失败。message 落 wiki_sources.parse_error。"""


# 不提取正文的 HTML 标签（script/style/title/meta 等）。
_SKIP_TAGS = {"script", "style", "title", "noscript", "head", "meta", "link"}


class _TagStripper(HTMLParser):
    """stdlib HTMLParser 提纯文本：忽略 script/style 等，取其余 data，压空白。

    MVP 级别 —— 真 boilerplate 抽取（readability-lxml/trafilatura）排期。
    对 webclip 已抓取的相对干净 HTML 够用；对杂乱页面质量一般。
    """

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._pieces: List[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag: str, attrs) -> None:  # noqa: ANN001
        if tag in _SKIP_TAGS:
            self._skip_depth += 1

    def handle_endtag(self, tag: str) -> None:
        if tag in _SKIP_TAGS and self._skip_depth > 0:
            self._skip_depth -= 1

    def handle_data(self, data: str) -> None:
        if self._skip_depth == 0:
            self._pieces.append(data)

    def text(self) -> str:
        raw = " ".join(self._pieces)
        return " ".join(raw.split())  # 压连续空白


def _looks_pdf(mime: str, filename: str) -> bool:
    return "pdf" in mime or filename.lower().endswith(".pdf")


def _looks_docx(mime: str, filename: str) -> bool:
    return "wordprocessingml" in mime or filename.lower().endswith(".docx")


def _looks_xlsx(mime: str, filename: str) -> bool:
    return "spreadsheetml" in mime or filename.lower().endswith(".xlsx")


def _looks_pptx(mime: str, filename: str) -> bool:
    return "presentationml" in mime or filename.lower().endswith(".pptx")


def _looks_epub(mime: str, filename: str) -> bool:
    return "epub" in mime or filename.lower().endswith(".epub")


def _looks_html(data: bytes, mime: str, filename: str) -> bool:
    if "html" in mime:
        return True
    name = filename.lower()
    if name.endswith((".html", ".htm", ".xhtml")):
        return True
    head = data[:512].lstrip().lower()
    return head.startswith(b"<!doctype html") or head.startswith(b"<html")


def count_pages(data: bytes, *, mime: str = "", filename: str = "") -> Optional[int]:
    """页数（云端解析按页计费用，client-docproc W4）。

    仅 PDF 有页概念（pypdf len(pages)）；其他格式返回 None（调用方
    自行兜底，如按 1 页或按字符）。解析失败同样 None —— 页数是计费
    元数据，绝不让它阻断主解析流程。
    """
    if not _looks_pdf(mime, filename):
        return None
    try:
        from pypdf import PdfReader
        return len(PdfReader(io.BytesIO(data)).pages)
    except Exception:  # noqa: BLE001 — 页数拿不到就当未知
        return None


def extract(data: bytes, *, mime: str = "", filename: str = "") -> str:
    """从文件字节提取纯文本。

    分派：PDF（mime/扩展名）→ DOCX/XLSX/PPTX/EPUB（mime/扩展名）→
    HTML（探测）→ 其余按 utf-8 decode（MD/TXT/code/JSON…）。
    空文件或无文本抛 ParseError。
    """
    if not data:
        raise ParseError("empty file")

    if _looks_pdf(mime, filename):
        return _extract_pdf(data)

    if _looks_docx(mime, filename):
        return _extract_docx(data)

    if _looks_xlsx(mime, filename):
        return _extract_xlsx(data)

    if _looks_pptx(mime, filename):
        return _extract_pptx(data)

    if _looks_epub(mime, filename):
        return _extract_epub(data)

    if _looks_html(data, mime, filename):
        return _extract_html(data)

    # MD/TXT/code/JSON/CSV/等 —— 直接 utf-8 decode（errors=replace 容错）
    try:
        text = data.decode("utf-8", errors="replace")
    except Exception as e:  # noqa: BLE001 — decode 几乎不抛（replace），兜底
        raise ParseError(f"utf-8 decode failed: {e}") from e
    if not text.strip():
        raise ParseError("decoded text empty")
    return text


def _extract_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise ParseError("pypdf not installed — add pypdf to worker deps") from e
    try:
        reader = PdfReader(io.BytesIO(data))
    except Exception as e:
        raise ParseError(f"open PDF failed: {e}") from e

    pieces: List[str] = []
    has_text = False
    for i, page in enumerate(reader.pages):
        try:
            text = (page.extract_text() or "").strip()
        except Exception as e:  # noqa: BLE001 — 单页失败不致命，占位继续
            pieces.append(f"[page {i + 1}: extraction failed: {e}]")
            continue
        if text:
            has_text = True
            pieces.append(f"--- page {i + 1} ---\n{text}")
        else:
            # 空文本页占位（修静默丢页：混合 PDF 用户须知道缺了哪页；
            # OCR 未启用的部署里这是扫描页的唯一信号）
            pieces.append(f"[page {i + 1}: 无文本层]")

    if not has_text:
        raise ParseError(
            "PDF contained no extractable text（扫描件？OCR 未启用或失败 —— "
            "启用 BIUMIND_WIKI_PARSE_OCR_ENABLED 走 MinerU）"
        )
    return "\n\n".join(pieces)


def _extract_html(data: bytes) -> str:
    try:
        html = data.decode("utf-8", errors="replace")
    except Exception as e:  # noqa: BLE001
        raise ParseError(f"html decode failed: {e}") from e
    stripper = _TagStripper()
    try:
        stripper.feed(html)
    except Exception as e:  # noqa: BLE001 — HTMLParser 对畸形 HTML 宽容
        raise ParseError(f"html parse failed: {e}") from e
    text = stripper.text().strip()
    if not text:
        raise ParseError("HTML stripped to empty (boilerplate-only or no body)")
    return text


def _strip_html_bytes(raw: bytes) -> str:
    """HTML 字节 → 纯文本（EPUB 章节复用；解析失败返回空串，由调用方兜底）。"""
    stripper = _TagStripper()
    try:
        stripper.feed(raw.decode("utf-8", errors="replace"))
    except Exception:  # noqa: BLE001 — 单章畸形不致命，跳过
        return ""
    return stripper.text().strip()


def _extract_docx(data: bytes) -> str:
    try:
        import mammoth
    except ImportError as e:
        raise ParseError("mammoth not installed — add mammoth to worker deps") from e
    try:
        result = mammoth.extract_raw_text(io.BytesIO(data))
    except Exception as e:
        raise ParseError(f"open DOCX failed: {e}") from e
    text = (result.value or "").strip()
    if not text:
        raise ParseError("DOCX contained no extractable text")
    return text


def _extract_xlsx(data: bytes) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise ParseError("openpyxl not installed — add openpyxl to worker deps") from e
    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as e:
        raise ParseError(f"open XLSX failed: {e}") from e
    try:
        sheets: List[str] = []
        for ws in wb.worksheets:
            rows: List[str] = []
            for row in ws.iter_rows(values_only=True):
                cells = [
                    str(c).strip()
                    for c in row
                    if c is not None and str(c).strip()
                ]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                sheets.append(f"--- sheet: {ws.title} ---\n" + "\n".join(rows))
    finally:
        wb.close()
    if not sheets:
        raise ParseError("XLSX contained no extractable text")
    return "\n\n".join(sheets)


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ParseError(
            "python-pptx not installed — add python-pptx to worker deps"
        ) from e
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:
        raise ParseError(f"open PPTX failed: {e}") from e
    slides: List[str] = []
    for i, slide in enumerate(prs.slides):
        texts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    texts.append(t)
        if texts:
            slides.append(f"--- slide {i + 1} ---\n" + "\n".join(texts))
    if not slides:
        raise ParseError("PPTX contained no extractable text")
    return "\n\n".join(slides)


def _extract_epub(data: bytes) -> str:
    try:
        import ebooklib
        from ebooklib import epub
    except ImportError as e:
        raise ParseError("ebooklib not installed — add ebooklib to worker deps") from e
    try:
        book = epub.read_epub(io.BytesIO(data))
    except Exception as e:
        raise ParseError(f"open EPUB failed: {e}") from e
    chapters: List[str] = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        if isinstance(item, epub.EpubNav):
            continue  # nav.xhtml 是目录 boilerplate，不是正文
        text = _strip_html_bytes(item.get_content())
        if text:
            chapters.append(f"--- {item.get_name()} ---\n{text}")
    if not chapters:
        raise ParseError("EPUB contained no extractable text")
    return "\n\n".join(chapters)
